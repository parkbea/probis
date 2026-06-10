import http.server, json, os, base64
from urllib import request as ureq
from urllib.parse import urlparse, parse_qs, quote

PORT   = 8080
ROOT   = os.path.dirname(os.path.abspath(__file__))
PUBLIC = os.path.join(ROOT, 'public')
DATA   = os.path.join(ROOT, 'data')

def op_call(api_path, base_url, api_key):
    url  = base_url.rstrip('/') + api_path
    auth = base64.b64encode(f'apikey:{api_key}'.encode()).decode()
    req  = ureq.Request(url, headers={'Authorization': f'Basic {auth}', 'Accept': 'application/json'})
    with ureq.urlopen(req, timeout=15) as resp:
        return resp.status, json.loads(resp.read())

def build_pptx(from_date, to_date, next_from, next_to, rows):
    """PPTX 생성 — 외부 라이브러리 없이 zipfile + XML만 사용"""
    import zipfile, io

    def ex(s):
        return str(s or '').replace('&','&amp;').replace('<','&lt;')\
                           .replace('>','&gt;').replace('"','&quot;')\
                           .replace("'",'&apos;')

    SW = 12192000   # slide width  EMU (≈13.33 in)
    SH =  6858000   # slide height EMU (=7.5  in)

    # ── 기본 XML 조각 ────────────────────────────────────────────

    def rpr(color='FFFFFF', sz=1200, bold=False):
        b = ' b="1"' if bold else ''
        return (f'<a:rPr lang="ko-KR" sz="{sz}"{b} dirty="0">'
                f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                f'<a:latin typeface="맑은 고딕"/>'
                f'</a:rPr>')

    def paragraphs(text, color='FFFFFF', sz=1200, bold=False, align='l'):
        out = ''
        for line in str(text or '').split('\n'):
            out += (f'<a:p><a:pPr algn="{align}"/>'
                    f'<a:r>{rpr(color,sz,bold)}<a:t>{ex(line)}</a:t></a:r>'
                    f'</a:p>')
        return out or '<a:p/>'

    def txbody(text, color='FFFFFF', sz=1200, bold=False, align='l', wrap='square'):
        return (f'<p:txBody><a:bodyPr wrap="{wrap}" rtlCol="0"/>'
                f'<a:lstStyle/>{paragraphs(text,color,sz,bold,align)}</p:txBody>')

    def shape(sid, name, ox, oy, cx, cy, fill=None, text='',
              sz=1200, bold=False, color='FFFFFF', align='l', wrap='square'):
        fill_xml = (f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
                    if fill else '<a:noFill/>')
        ln_xml   = '' if fill else ''
        return (f'<p:sp>'
                f'<p:nvSpPr>'
                f'<p:cNvPr id="{sid}" name="{ex(name)}"/>'
                f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
                f'<p:nvPr/></p:nvSpPr>'
                f'<p:spPr>'
                f'<a:xfrm><a:off x="{ox}" y="{oy}"/>'
                f'<a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
                f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
                f'{fill_xml}<a:ln><a:noFill/></a:ln>'
                f'</p:spPr>'
                f'{txbody(text,color,sz,bold,align,wrap)}'
                f'</p:sp>')

    # ── 테이블 ───────────────────────────────────────────────────

    def tc_xml(text, fill=None, color='1E293B', sz=850, bold=False, align='l'):
        fill_xml = (f'<a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>'
                    if fill else '<a:noFill/>')
        bdr = ('<a:lnL w="9525"><a:solidFill><a:srgbClr val="CBD5E1"/></a:solidFill></a:lnL>'
               '<a:lnR w="9525"><a:solidFill><a:srgbClr val="CBD5E1"/></a:solidFill></a:lnR>'
               '<a:lnT w="9525"><a:solidFill><a:srgbClr val="CBD5E1"/></a:solidFill></a:lnT>'
               '<a:lnB w="9525"><a:solidFill><a:srgbClr val="CBD5E1"/></a:solidFill></a:lnB>')
        body = (f'<a:txBody><a:bodyPr wrap="square"/><a:lstStyle/>'
                f'{paragraphs(text,color,sz,bold,align)}</a:txBody>')
        pr   = f'<a:tcPr marL="91440" marR="91440" marT="45720" marB="45720">{fill_xml}{bdr}</a:tcPr>'
        return f'<a:tc>{body}{pr}</a:tc>'

    def table_xml(ox, oy, cx, cy, col_ws, headers, data_rows):
        grid = ''.join(f'<a:gridCol w="{w}"/>' for w in col_ws)
        hdr_h  = 411480
        row_h  = max(330000, min(914400, (cy - hdr_h) // max(len(data_rows), 1)))

        hdr_row = '<a:tr h="{}">{}</a:tr>'.format(
            hdr_h,
            ''.join(tc_xml(h, fill='1E293B', color='FFFFFF', sz=1000, bold=True)
                    for h in headers))

        data = ''
        for ri, row_vals in enumerate(data_rows):
            fill = 'F0FDFA' if ri % 2 == 0 else 'FFFFFF'
            cells = ''.join(
                tc_xml(v, fill=fill, bold=(ci == 0))
                for ci, v in enumerate(row_vals))
            data += f'<a:tr h="{row_h}">{cells}</a:tr>'

        return (f'<p:graphicFrame>'
                f'<p:nvGraphicFramePr>'
                f'<p:cNvPr id="10" name="Table"/>'
                f'<p:cNvGraphicFramePr>'
                f'<a:graphicFrameLocks noGrp="1"/>'
                f'</p:cNvGraphicFramePr><p:nvPr/></p:nvGraphicFramePr>'
                f'<p:xfrm><a:off x="{ox}" y="{oy}"/>'
                f'<a:ext cx="{cx}" cy="{cy}"/></p:xfrm>'
                f'<a:graphic>'
                f'<a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">'
                f'<a:tbl><a:tblPr/>'
                f'<a:tblGrid>{grid}</a:tblGrid>'
                f'{hdr_row}{data}'
                f'</a:tbl></a:graphicData></a:graphic>'
                f'</p:graphicFrame>')

    # ── 슬라이드 XML 래퍼 ────────────────────────────────────────

    def slide_wrap(content):
        return (f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                f'<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
                f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                f'<p:cSld><p:spTree>'
                f'<p:nvGrpSpPr><p:cNvPr id="1" name=""/>'
                f'<p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
                f'<p:grpSpPr><a:xfrm>'
                f'<a:off x="0" y="0"/><a:ext cx="{SW}" cy="{SH}"/>'
                f'<a:chOff x="0" y="0"/><a:chExt cx="{SW}" cy="{SH}"/>'
                f'</a:xfrm></p:grpSpPr>'
                f'{content}'
                f'</p:spTree></p:cSld>'
                f'<p:clrMapOvr><a:masterClr/></p:clrMapOvr>'
                f'</p:sld>')

    def mmdd(s):
        p = str(s or '').split('-')
        return f'{p[1]}/{p[2]}' if len(p) >= 3 else s
    def ymd(s):
        return str(s or '').replace('-', '/')

    # ── 슬라이드 1: 테이블 (헤더 텍스트 없음) ───────────────────
    margin = 274638
    tbl_x  = margin
    tbl_y  = 182880
    tbl_cx = SW - margin * 2
    tbl_cy = SH - tbl_y - 182880

    col_pcts = [145, 220, 220, 105, 80, 80, 150]  # /1000
    col_ws   = [tbl_cx * p // 1000 for p in col_pcts]
    col_ws[-1] += tbl_cx - sum(col_ws)

    hdrs = ['프로젝트명',
            f'금주 내용\n({mmdd(from_date)} ~ {mmdd(to_date)})',
            f'차주 진행 내용\n({mmdd(next_from)} ~ {mmdd(next_to)})',
            '담당자','시작일','종료일','비고']
    data = [[r.get('name',''), r.get('thisWeek',''), r.get('nextWeek',''),
             r.get('assignees',''), ymd(r.get('startDate','')),
             ymd(r.get('endDate','')), r.get('note','')]
            for r in rows] if rows else [['(데이터 없음)']+['']*6]

    s1 = table_xml(tbl_x, tbl_y, tbl_cx, tbl_cy, col_ws, hdrs, data)

    # ── 정적 XML ─────────────────────────────────────────────────
    slide_rels = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                  '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                  '<Relationship Id="rId1" '
                  'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" '
                  'Target="../slideLayouts/slideLayout1.xml"/>'
                  '</Relationships>')

    files = {
        '[Content_Types].xml': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml"  ContentType="application/xml"/>'
            '<Override PartName="/ppt/presentation.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>'
            '<Override PartName="/ppt/slides/slide1.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
            '<Override PartName="/ppt/slideLayouts/slideLayout1.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>'
            '<Override PartName="/ppt/slideMasters/slideMaster1.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>'
            '<Override PartName="/ppt/theme/theme1.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
            '<Override PartName="/ppt/presProps.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.presProps+xml"/>'
            '<Override PartName="/ppt/tableStyles.xml"'
            ' ContentType="application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml"/>'
            '</Types>'),

        '_rels/.rels': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"'
            ' Target="ppt/presentation.xml"/>'
            '</Relationships>'),

        'ppt/presentation.xml': (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            f'<p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>'
            f'<p:sldIdLst>'
            f'<p:sldId id="256" r:id="rId2"/>'
            f'</p:sldIdLst>'
            f'<p:sldSz cx="{SW}" cy="{SH}" type="custom"/>'
            f'<p:notesSz cx="6858000" cy="9144000"/>'
            f'</p:presentation>'),

        'ppt/_rels/presentation.xml.rels': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"'
            ' Target="slideMasters/slideMaster1.xml"/>'
            '<Relationship Id="rId2"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"'
            ' Target="slides/slide1.xml"/>'
            '<Relationship Id="rId3"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/presProps"'
            ' Target="presProps.xml"/>'
            '<Relationship Id="rId4"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/tableStyles"'
            ' Target="tableStyles.xml"/>'
            '</Relationships>'),

        'ppt/theme/theme1.xml': (
            '<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="koFlow">'
            '<a:themeElements>'
            '<a:clrScheme name="koFlow">'
            '<a:dk1><a:srgbClr val="1E293B"/></a:dk1>'
            '<a:lt1><a:srgbClr val="FFFFFF"/></a:lt1>'
            '<a:dk2><a:srgbClr val="334155"/></a:dk2>'
            '<a:lt2><a:srgbClr val="F8FAFC"/></a:lt2>'
            '<a:accent1><a:srgbClr val="6366F1"/></a:accent1>'
            '<a:accent2><a:srgbClr val="0D9488"/></a:accent2>'
            '<a:accent3><a:srgbClr val="3B82F6"/></a:accent3>'
            '<a:accent4><a:srgbClr val="F59E0B"/></a:accent4>'
            '<a:accent5><a:srgbClr val="10B981"/></a:accent5>'
            '<a:accent6><a:srgbClr val="EF4444"/></a:accent6>'
            '<a:hlink><a:srgbClr val="6366F1"/></a:hlink>'
            '<a:folHlink><a:srgbClr val="4338CA"/></a:folHlink>'
            '</a:clrScheme>'
            '<a:fontScheme name="koFlow">'
            '<a:majorFont><a:latin typeface="맑은 고딕"/><a:ea typeface=""/><a:cs typeface=""/></a:majorFont>'
            '<a:minorFont><a:latin typeface="맑은 고딕"/><a:ea typeface=""/><a:cs typeface=""/></a:minorFont>'
            '</a:fontScheme>'
            '<a:fmtScheme name="koFlow">'
            '<a:fillStyleLst>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '</a:fillStyleLst>'
            '<a:lnStyleLst>'
            '<a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
            '<a:ln w="12700"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
            '<a:ln w="19050"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>'
            '</a:lnStyleLst>'
            '<a:effectStyleLst>'
            '<a:effectStyle><a:effectLst/></a:effectStyle>'
            '<a:effectStyle><a:effectLst/></a:effectStyle>'
            '<a:effectStyle><a:effectLst/></a:effectStyle>'
            '</a:effectStyleLst>'
            '<a:bgFillStyleLst>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '<a:solidFill><a:schemeClr val="phClr"/></a:solidFill>'
            '</a:bgFillStyleLst>'
            '</a:fmtScheme>'
            '</a:themeElements>'
            '</a:theme>'),

        'ppt/slideMasters/slideMaster1.xml': (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            f'<p:cSld><p:spTree>'
            f'<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{SW}" cy="{SH}"/>'
            f'<a:chOff x="0" y="0"/><a:chExt cx="{SW}" cy="{SH}"/></a:xfrm></p:grpSpPr>'
            f'</p:spTree></p:cSld>'
            f'<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2"'
            f' accent1="accent1" accent2="accent2" accent3="accent3"'
            f' accent4="accent4" accent5="accent5" accent6="accent6"'
            f' hlink="hlink" folHlink="folHlink"/>'
            f'<p:sldLayoutIdLst>'
            f'<p:sldLayoutId id="2147483649" r:id="rId2"/>'
            f'</p:sldLayoutIdLst>'
            f'<p:txStyles>'
            f'<p:titleStyle><a:lvl1pPr><a:defRPr lang="ko-KR"/></a:lvl1pPr></p:titleStyle>'
            f'<p:bodyStyle><a:lvl1pPr><a:defRPr lang="ko-KR"/></a:lvl1pPr></p:bodyStyle>'
            f'<p:otherStyle><a:defPPr><a:defRPr lang="ko-KR"/></a:defPPr></p:otherStyle>'
            f'</p:txStyles>'
            f'</p:sldMaster>'),

        'ppt/slideMasters/_rels/slideMaster1.xml.rels': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"'
            ' Target="../theme/theme1.xml"/>'
            '<Relationship Id="rId2"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"'
            ' Target="../slideLayouts/slideLayout1.xml"/>'
            '</Relationships>'),

        'ppt/slideLayouts/slideLayout1.xml': (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            f' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"'
            f' type="blank" preserve="1">'
            f'<p:cSld name="Blank"><p:spTree>'
            f'<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
            f'<p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="{SW}" cy="{SH}"/>'
            f'<a:chOff x="0" y="0"/><a:chExt cx="{SW}" cy="{SH}"/></a:xfrm></p:grpSpPr>'
            f'</p:spTree></p:cSld>'
            f'<p:clrMapOvr><a:masterClr/></p:clrMapOvr>'
            f'</p:sldLayout>'),

        'ppt/slideLayouts/_rels/slideLayout1.xml.rels': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1"'
            ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"'
            ' Target="../slideMasters/slideMaster1.xml"/>'
            '</Relationships>'),

        'ppt/slides/slide1.xml':            slide_wrap(s1),
        'ppt/slides/_rels/slide1.xml.rels': slide_rels,

        'ppt/presProps.xml': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<p:presProps xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:showPr/></p:presProps>'),

        'ppt/tableStyles.xml': (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<a:tblStyleLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
            ' def="{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}"/>'),
    }

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zf:
        for name, content in files.items():
            zf.writestr(name, content.encode('utf-8'))
    return buf.getvalue()


def build_pptx_advanced(from_date, to_date, next_from, next_to, rows):
    """python-pptx 기반 고품질 PPTX (pip install python-pptx 또는 install_packages.bat 필요)"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    W, H = Inches(13.33), Inches(7.5)
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]

    def mmdd(s):
        p = str(s or '').split('-')
        return f'{p[1]}/{p[2]}' if len(p) >= 3 else s
    def ymd(s):
        return str(s or '').replace('-', '/')

    # ── 슬라이드 1: 테이블 (헤더 텍스트 없음) ───────────────────
    s1 = prs.slides.add_slide(blank)

    disp_rows = rows if rows else [{}]
    tbl = s1.shapes.add_table(
        len(disp_rows)+1, 7, Inches(0.3), Inches(0.2), Inches(12.73), Inches(7.1)
    ).table
    for i, w in enumerate([Inches(1.8), Inches(2.83), Inches(2.83),
                            Inches(1.3), Inches(0.95), Inches(0.95), Inches(1.07)]):
        tbl.columns[i].width = w

    def cell(c, text, bg=None, fg=RGBColor(0x1E,0x29,0x3B), bold=False, sz=9):
        c.text = ''
        tf = c.text_frame; tf.word_wrap = True
        tf.margin_left = tf.margin_right = Emu(60000)
        tf.margin_top  = tf.margin_bottom = Emu(50000)
        for idx, line in enumerate(str(text or '').split('\n')):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = line
            run.font.size = Pt(sz); run.font.bold = bold
            run.font.color.rgb = fg
        if bg: c.fill.solid(); c.fill.fore_color.rgb = bg

    col_headers = [
        '프로젝트명',
        f'금주 내용\n({mmdd(from_date)} ~ {mmdd(to_date)})',
        f'차주 진행 내용\n({mmdd(next_from)} ~ {mmdd(next_to)})',
        '담당자','시작일','종료일','비고',
    ]
    for ci, h in enumerate(col_headers):
        cell(tbl.cell(0,ci), h, bg=RGBColor(0x1E,0x29,0x3B),
             fg=RGBColor(0xFF,0xFF,0xFF), bold=True, sz=10)
    tbl.rows[0].height = Inches(0.45)

    rh = int(max(Inches(0.38), min(Inches(1.1),
             (Inches(7.1)-Inches(0.45)) // max(len(disp_rows),1))))
    for ri, row in enumerate(disp_rows):
        bg = RGBColor(0xF0,0xFD,0xFA) if ri%2==0 else RGBColor(0xFF,0xFF,0xFF)
        for ci, val in enumerate([row.get('name',''), row.get('thisWeek',''),
                                   row.get('nextWeek',''), row.get('assignees',''),
                                   ymd(row.get('startDate','')), ymd(row.get('endDate','')),
                                   row.get('note','')]):
            cell(tbl.cell(ri+1,ci), val, bg=bg, bold=(ci==0))
        tbl.rows[ri+1].height = rh

    buf = io.BytesIO(); prs.save(buf)
    return buf.getvalue()


class Handler(http.server.SimpleHTTPRequestHandler):
    def __init__(self, *args, **kwargs):
        super().__init__(*args, directory=PUBLIC, **kwargs)

    def do_GET(self):
        parsed = urlparse(self.path)

        # ── OpenProject 동기화 ──────────────────
        if parsed.path == '/op-sync':
            self._handle_op_sync(); return

        # ── OpenProject 프록시 ──────────────────
        if parsed.path == '/op-proxy':
            qs          = parse_qs(parsed.query)
            target_path = qs.get('path', [''])[0]
            if not target_path:
                self.send_error(400, 'missing path'); return
            try:
                with open(os.path.join(DATA, 'config.json'), 'r', encoding='utf-8') as f:
                    cfg = json.load(f)
            except Exception:
                self.send_error(500, 'data/config.json 없음'); return
            op       = cfg.get('openproject', {})
            base_url = op.get('baseUrl', '').rstrip('/')
            api_key  = op.get('apiKey', '')
            if not base_url:
                self.send_error(400, 'OpenProject URL 미설정'); return
            try:
                status, body = op_call(target_path, base_url, api_key)
                data = json.dumps(body).encode()
                self.send_response(status)
                self.send_header('Content-Type', 'application/json')
                self.end_headers()
                self.wfile.write(data)
            except Exception as e:
                self.send_error(502, str(e))
            return

        # ── data/ JSON 읽기 ─────────────────────
        if parsed.path.startswith('/data/'):
            fname = os.path.basename(parsed.path)
            if not fname.endswith('.json'):
                self.send_error(403); return
            fpath = os.path.join(DATA, fname)
            try:
                with open(fpath, 'rb') as f:
                    data = f.read()
                self.send_response(200)
                self.send_header('Content-Type', 'application/json; charset=utf-8')
                self.end_headers()
                self.wfile.write(data)
            except FileNotFoundError:
                self.send_error(404)
            return

        super().do_GET()

    def _handle_op_sync(self):
        try:
            with open(os.path.join(DATA, 'config.json'),  'r', encoding='utf-8') as f:
                cfg = json.load(f)
            with open(os.path.join(DATA, 'members.json'), 'r', encoding='utf-8') as f:
                members_data = json.load(f)
        except Exception as e:
            self._json_response(500, {'error': f'파일 읽기 실패: {e}'}); return

        op       = cfg.get('openproject', {})
        base_url = op.get('baseUrl', '').rstrip('/')
        api_key  = op.get('apiKey', '')
        if not base_url or not api_key:
            self._json_response(400, {'error': 'OpenProject 설정(URL/API 키)이 필요합니다'}); return

        # opUserId가 등록된 팀원만 사용 (opId 쿼리 지정 시 해당 직원만 갱신)
        from urllib.parse import urlparse, parse_qs
        only_op_id = (parse_qs(urlparse(self.path).query).get('opId') or [None])[0]
        members  = members_data.get('members', [])
        op_users = [
            {'name': m.get('name', ''), 'opId': str(m['opUserId'])}
            for m in members
            if m.get('opUserId') and (not only_op_id or str(m['opUserId']) == str(only_op_id))
        ]

        if not op_users:
            self._json_response(400, {
                'error': '팀원 관리에서 OpenProject 사용자 ID를 먼저 입력해주세요.',
            }); return

        # 2. 사용자에게 '할당된(assignee)' 작업이 있는 프로젝트만 조회
        user_ids = [u['opId'] for u in op_users]
        wp_flt   = quote(json.dumps([{'assignee': {'operator': '=', 'values': user_ids}}]))
        try:
            status, wp_body = op_call(f'/api/v3/work_packages?filters={wp_flt}&pageSize=200', base_url, api_key)
        except Exception as e:
            self._json_response(502, {'error': f'OpenProject 작업 조회 실패: {e}'}); return

        # 프로젝트별로 '우리 직원 중 담당자(opId)' 모으기
        proj_assignees = {}
        for wp in wp_body.get('_embedded', {}).get('elements', []):
            links = wp.get('_links', {})
            phref = (links.get('project') or {}).get('href', '')
            ahref = (links.get('assignee') or {}).get('href', '')
            pid   = phref.rstrip('/').split('/')[-1] if phref else ''
            aid   = ahref.rstrip('/').split('/')[-1] if ahref else ''
            if not pid:
                continue
            proj_assignees.setdefault(pid, set())
            if aid:
                proj_assignees[pid].add(aid)
        proj_ids = list(proj_assignees.keys())

        if not proj_ids:
            self._json_response(200, {
                'projects': [], 'opUsers': op_users,
                'message': '할당된 프로젝트가 없습니다.',
            }); return

        # 그 프로젝트들의 상세 정보 조회
        pj_flt = quote(json.dumps([{'id': {'operator': '=', 'values': proj_ids}}]))
        try:
            status, body = op_call(f'/api/v3/projects?filters={pj_flt}&pageSize=200', base_url, api_key)
        except Exception as e:
            self._json_response(502, {'error': f'OpenProject 프로젝트 조회 실패: {e}'}); return

        # 상태가 '마침(finished)'인 프로젝트는 제외
        def _is_finished(p):
            st    = p.get('_links', {}).get('status') or {}
            href  = (st.get('href') or '').lower()
            title = st.get('title') or ''
            return href.endswith('/finished') or 'finish' in title.lower() or '마침' in title

        op_projects = [
            p for p in body.get('_embedded', {}).get('elements', [])
            if not _is_finished(p)
        ]

        # 3. koFlow 포맷으로 변환
        from datetime import datetime, timezone
        now_iso = datetime.now(timezone.utc).isoformat()
        mapped = [{
            'id':           'op_' + str(p['id']),
            'opUserIds':    list(proj_assignees.get(str(p['id']), set())),
            'type':         '실행중인 프로젝트',
            'name':         p.get('name', ''),
            'startDate':    p.get('startDate') or '',
            'endDate':      p.get('dueDate') or '',
            'status':       '완료' if 'finish' in (p.get('status', {}).get('name') or '').lower() else '진행중',
            'effort':       0,
            'effortUnit':   'MM',
            'opEpicUrl':    f"{base_url}/projects/{p.get('identifier', '')}",
            'opEffortUrl':  '',
            'opQaUrl':      '',
            'emailContent': '',
            'assignments':  [],
            'createdAt':    p.get('createdAt') or now_iso,
            'updatedAt':    now_iso,
        } for p in op_projects]

        self._json_response(200, {
            'projects': mapped,
            'opUsers':  op_users,
            'message':  f'{len(op_users)}명의 팀원이 속한 프로젝트 {len(mapped)}개를 찾았습니다.',
        })

    def _handle_weekly_report(self, body_bytes):
        try:
            data = json.loads(body_bytes)
        except Exception:
            self._json_response(400, {'error': '잘못된 요청 데이터'}); return

        from_date = data.get('from',     '')
        to_date   = data.get('to',       '')
        next_from = data.get('nextFrom', '')
        next_to   = data.get('nextTo',   '')
        rows      = data.get('rows',     [])
        try:
            ppt_bytes = build_pptx_advanced(from_date, to_date, next_from, next_to, rows)
        except ImportError:
            ppt_bytes = build_pptx(from_date, to_date, next_from, next_to, rows)
        self.send_response(200)
        self.send_header('Content-Type',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        self.send_header('Content-Disposition', 'attachment; filename="weekly_report.pptx"')
        self.send_header('Content-Length', str(len(ppt_bytes)))
        self.end_headers()
        self.wfile.write(ppt_bytes)

    def _json_response(self, status, data):
        body = json.dumps(data, ensure_ascii=False).encode('utf-8')
        self.send_response(status)
        self.send_header('Content-Type', 'application/json; charset=utf-8')
        self.end_headers()
        self.wfile.write(body)

    def do_POST(self):
        parsed = urlparse(self.path)

        # ── 주간보고 PPT 생성 ───────────────────────
        if parsed.path == '/weekly-report':
            length = int(self.headers.get('Content-Length', 0))
            body   = self.rfile.read(length)
            self._handle_weekly_report(body)
            return

        if not parsed.path.startswith('/data/'):
            self.send_error(403); return
        fname = os.path.basename(parsed.path)
        if not fname.endswith('.json'):
            self.send_error(403); return
        fpath = os.path.join(DATA, fname)
        length = int(self.headers.get('Content-Length', 0))
        body   = self.rfile.read(length)
        try:
            json.loads(body)
            with open(fpath, 'wb') as f:
                f.write(body)
            self.send_response(200)
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(b'{"ok":true}')
        except Exception as e:
            self.send_error(400, str(e))

    def log_message(self, fmt, *args):
        pass

if __name__ == '__main__':
    os.chdir(ROOT)
    with http.server.HTTPServer(('127.0.0.1', PORT), Handler) as s:
        print(f'koFlow 서버 실행 중: http://localhost:{PORT}')
        s.serve_forever()
